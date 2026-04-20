"""
RAG document store skill.

Ingests uploaded files (PDF, TXT, MD, DOCX, XLSX, PPTX) into a Qdrant
collection keyed by research session ID. Returns stored point IDs so the
pipeline can retrieve context for each LLM's research prompt.

Dependencies: qdrant-client, openai (embeddings), pypdf,
              python-docx, openpyxl, python-pptx
"""

import io
import logging
import os
import uuid
from typing import BinaryIO

logger = logging.getLogger(__name__)

_COLLECTION = "market_research_rag"
_EMBED_MODEL = "text-embedding-3-small"
_CHUNK_SIZE   = 800   # characters per chunk
_CHUNK_OVERLAP = 100

# ── Qdrant client (lazy) ───────────────────────────────────────────────────────

_qdrant_client = None

def _get_qdrant():
    global _qdrant_client
    if _qdrant_client is None:
        qdrant_url = os.environ.get("QDRANT_URL")
        if not qdrant_url:
            logger.warning("rag_store: QDRANT_URL not set — RAG disabled")
            return None
        from qdrant_client import QdrantClient
        _qdrant_client = QdrantClient(
            url=qdrant_url,
            api_key=os.environ.get("QDRANT_API_KEY"),
        )
        _ensure_collection(_qdrant_client)
    return _qdrant_client


def _ensure_collection(client) -> None:
    from qdrant_client.models import Distance, VectorParams
    existing = {c.name for c in client.get_collections().collections}
    if _COLLECTION not in existing:
        client.create_collection(
            collection_name=_COLLECTION,
            vectors_config=VectorParams(size=1536, distance=Distance.COSINE),
        )
        logger.info("rag_store: created Qdrant collection '%s'", _COLLECTION)


# ── Text extraction ────────────────────────────────────────────────────────────

def _extract_text(filename: str, data: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if ext in ("docx", "doc"):
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if ext in ("xlsx", "xls"):
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
        return "\n".join(lines)
    if ext in ("pptx", "ppt"):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)
    # plain text fallback (txt, md, csv, etc.)
    return data.decode("utf-8", errors="replace")


def _chunk(text: str) -> list[str]:
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = start + _CHUNK_SIZE
        chunks.append(text[start:end])
        start += _CHUNK_SIZE - _CHUNK_OVERLAP
    return [c.strip() for c in chunks if c.strip()]


# ── Embedding ──────────────────────────────────────────────────────────────────

def _embed(texts: list[str]) -> list[list[float]]:
    from openai import OpenAI
    client = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
    resp = client.embeddings.create(model=_EMBED_MODEL, input=texts)
    return [item.embedding for item in resp.data]


# ── Public API ─────────────────────────────────────────────────────────────────

def ingest_document(filename: str, data: bytes, session_id: str) -> list[str]:
    """
    Parse, chunk, embed, and upsert a document into Qdrant.

    Returns list of point ID strings stored for this document.
    """
    client = _get_qdrant()
    if client is None:
        logger.warning("rag_store: skipping ingest — Qdrant unavailable")
        return []

    text = _extract_text(filename, data)
    if not text.strip():
        logger.warning("rag_store: %s produced no text", filename)
        return []

    chunks = _chunk(text)
    logger.info("rag_store: %s → %d chunks", filename, len(chunks))

    embeddings = _embed(chunks)

    from qdrant_client.models import PointStruct
    point_ids = [str(uuid.uuid4()) for _ in chunks]
    points = [
        PointStruct(
            id=pid,
            vector=emb,
            payload={
                "session_id": session_id,
                "filename":   filename,
                "chunk_idx":  i,
                "text":       chunk,
            },
        )
        for i, (pid, emb, chunk) in enumerate(zip(point_ids, embeddings, chunks))
    ]
    client.upsert(collection_name=_COLLECTION, points=points)
    logger.info("rag_store: upserted %d points for session %s", len(points), session_id)
    return point_ids


def retrieve_context(query: str, session_id: str, top_k: int = 5) -> str:
    """
    Retrieve the most relevant chunks for a query within a session.

    Returns a single string to be appended to LLM prompts.
    """
    try:
        client = _get_qdrant()
        if client is None:
            return ""
        query_vec = _embed([query])[0]
        hits = client.search(
            collection_name=_COLLECTION,
            query_vector=query_vec,
            limit=top_k,
            query_filter={
                "must": [{"key": "session_id", "match": {"value": session_id}}]
            },
        )
        if not hits:
            return ""
        snippets = [h.payload["text"] for h in hits]
        return "\n\n---\n\n".join(snippets)
    except Exception as exc:
        logger.error("rag_store: retrieval failed — %s", exc)
        return ""


def delete_session_docs(session_id: str) -> None:
    """Remove all vectors for a completed research session."""
    try:
        client = _get_qdrant()
        client.delete(
            collection_name=_COLLECTION,
            points_selector={"filter": {
                "must": [{"key": "session_id", "match": {"value": session_id}}]
            }},
        )
    except Exception as exc:
        logger.error("rag_store: delete failed for session %s — %s", session_id, exc)
